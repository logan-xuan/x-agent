from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

def create_spring_story_ppt():
    """Create a presentation about the story of spring."""
    prs = Presentation()
    
    # Set theme colors - fresh spring greens and blues
    prs.slide_master.background.fill.solid()
    prs.slide_master.background.fill.fore_color.rgb = RGBColor(245, 255, 245)  # Light green background
    
    # Define colors for consistent styling
    title_color = RGBColor(30, 80, 40)  # Dark green for titles
    text_color = RGBColor(50, 90, 60)   # Medium green for text
    accent_color = RGBColor(255, 215, 100)  # Gold for accents
    
    # Slide 1: Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "春天的故事"
    subtitle.text = "生机勃勃的季节 - 感受春天的美好"
    
    # Style title and subtitle
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.color.rgb = title_color
    title.text_frame.paragraphs[0].font.bold = True
    subtitle.text_frame.paragraphs[0].font.size = Pt(24)
    subtitle.text_frame.paragraphs[0].font.color.rgb = text_color
    
    # Slide 2: Spring's Arrival
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "春天的序曲"
    title.text_frame.paragraphs[0].font.color.rgb = title_color
    title.text_frame.paragraphs[0].font.size = Pt(36)
    
    content.text = ("冬去春来，大地渐渐苏醒\n"
                   "• 温暖的阳光洒向大地\n"
                   "• 雪花融化成清澈的溪流\n"
                   "• 鸟儿开始欢快地歌唱")
    
    for paragraph in content.text_frame.paragraphs[1:]:
        paragraph.font.size = Pt(18)
        paragraph.font.color.rgb = text_color
        paragraph.level = 0
    
    # Slide 3: Awakening Earth
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "大地苏醒"
    title.text_frame.paragraphs[0].font.color.rgb = title_color
    title.text_frame.paragraphs[0].font.size = Pt(36)
    
    content.text = ("植物的复苏与生长\n"
                   "• 嫩绿的新芽破土而出\n"
                   "• 枝条上长出翠绿的叶片\n"
                   "• 草地重新披上绿色外衣")
    
    for paragraph in content.text_frame.paragraphs[1:]:
        paragraph.font.size = Pt(18)
        paragraph.font.color.rgb = text_color
        paragraph.level = 0
    
    # Slide 4: Animal World
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "动物归来"
    title.text_frame.paragraphs[0].font.color.rgb = title_color
    title.text_frame.paragraphs[0].font.size = Pt(36)
    
    content.text = ("春天里的动物世界\n"
                   "• 迁徙的鸟儿返回故乡\n"
                   "• 昆虫开始活跃起来\n"
                   "• 小动物们结束冬眠")
    
    for paragraph in content.text_frame.paragraphs[1:]:
        paragraph.font.size = Pt(18)
        paragraph.font.color.rgb = text_color
        paragraph.level = 0
    
    # Slide 5: Spring Rain
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "春雨滋润"
    title.text_frame.paragraphs[0].font.color.rgb = title_color
    title.text_frame.paragraphs[0].font.size = Pt(36)
    
    content.text = ("雨水对生命的意义\n"
                   "• 细雨滋润干涸的土地\n"
                   "• 为植物提供必需的水分\n"
                   "• 带来清新怡人的空气")
    
    for paragraph in content.text_frame.paragraphs[1:]:
        paragraph.font.size = Pt(18)
        paragraph.font.color.rgb = text_color
        paragraph.level = 0
    
    # Slide 6: Blooming Flowers
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "花开满园"
    title.text_frame.paragraphs[0].font.color.rgb = title_color
    title.text_frame.paragraphs[0].font.size = Pt(36)
    
    content.text = ("各种春花竞相绽放\n"
                   "• 桃花粉嫩如霞\n"
                   "• 樱花洁白似雪\n"
                   "• 油菜花金黄耀眼")
    
    for paragraph in content.text_frame.paragraphs[1:]:
        paragraph.font.size = Pt(18)
        paragraph.font.color.rgb = text_color
        paragraph.level = 0
    
    # Slide 7: People Welcoming Spring
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "人们迎接春天"
    title.text_frame.paragraphs[0].font.color.rgb = title_color
    title.text_frame.paragraphs[0].font.size = Pt(36)
    
    content.text = ("春季活动与习俗\n"
                   "• 踏青赏花游春\n"
                   "• 放风筝感受春风\n"
                   "• 种植花草树木")
    
    for paragraph in content.text_frame.paragraphs[1:]:
        paragraph.font.size = Pt(18)
        paragraph.font.color.rgb = text_color
        paragraph.level = 0
    
    # Slide 8: Spring Colors
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "春天的色彩"
    title.text_frame.paragraphs[0].font.color.rgb = title_color
    title.text_frame.paragraphs[0].font.size = Pt(36)
    
    content.text = ("自然界的缤纷色彩\n"
                   "• 新绿：生命的颜色\n"
                   "• 粉红：桃花的颜色\n"
                   "• 淡蓝：晴空的颜色")
    
    for paragraph in content.text_frame.paragraphs[1:]:
        paragraph.font.size = Pt(18)
        paragraph.font.color.rgb = text_color
        paragraph.level = 0
    
    # Slide 9: Symbolism of Spring
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "春天的寓意"
    title.text_frame.paragraphs[0].font.color.rgb = title_color
    title.text_frame.paragraphs[0].font.size = Pt(36)
    
    content.text = ("新生与希望的象征\n"
                   "• 万物复苏的季节\n"
                   "• 充满活力与生机\n"
                   "• 寓意美好未来")
    
    for paragraph in content.text_frame.paragraphs[1:]:
        paragraph.font.size = Pt(18)
        paragraph.font.color.rgb = text_color
        paragraph.level = 0
    
    # Slide 10: Conclusion
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "感受春天的美好"
    title.text_frame.paragraphs[0].font.color.rgb = title_color
    title.text_frame.paragraphs[0].font.size = Pt(36)
    
    content.text = ("春天的故事永不停歇\n"
                   "• 让我们拥抱春天\n"
                   "• 感受生命的奇迹\n"
                   "• 享受大自然的馈赠")
    
    for paragraph in content.text_frame.paragraphs[1:]:
        paragraph.font.size = Pt(18)
        paragraph.font.color.rgb = text_color
        paragraph.level = 0
    
    # CRITICAL: Save to presentations/ directory, NOT workspace root!
    output_dir = 'presentations'
    os.makedirs(output_dir, exist_ok=True)  # Create directory if needed
    
    output_path = os.path.join(output_dir, 'spring_story.pptx')
    prs.save(output_path)
    
    print(f"✅ PPTX saved to: {output_path}")
    return output_path

if __name__ == "__main__":
    create_spring_story_ppt()