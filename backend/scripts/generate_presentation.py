from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import re

def create_presentation_from_markdown(md_content):
    # Create presentation object
    prs = Presentation()
    
    # Set slide dimensions (standard 4:3 aspect ratio)
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # Define styles
    title_font_size = Pt(44)
    heading_font_size = Pt(36)
    content_font_size = Pt(24)
    
    # Parse markdown content
    lines = md_content.split('\n')
    
    # Process each line
    i = 0
    while i < len(lines):
        line = lines[i].strip()
        
        if line.startswith('# '):  # Main title
            title_slide = prs.slides.add_slide(prs.slide_layouts[0])  # Title slide layout
            title = title_slide.shapes.title
            subtitle = title_slide.placeholders[1]
            
            title.text = line[2:]  # Remove '# '
            title.text_frame.paragraphs[0].font.size = title_font_size
            title.text_frame.paragraphs[0].font.bold = True
            
            subtitle.text = "2026年AI发展趋势深度研究报告\n趋势分析与展望"
            subtitle.text_frame.paragraphs[0].font.size = Pt(28)
            
            i += 1
        elif line.startswith('## '):  # Section headings
            section_title = line[3:]  # Remove '## '
            
            # Add section slide
            blank_slide_layout = prs.slide_layouts[6]  # Blank slide
            section_slide = prs.slides.add_slide(blank_slide_layout)
            
            # Add title box
            title_box = section_slide.shapes.add_textbox(Inches(1), Inches(1), Inches(11.33), Inches(1.5))
            title_tf = title_box.text_frame
            title_tf.text = section_title
            title_tf.paragraphs[0].font.size = heading_font_size
            title_tf.paragraphs[0].font.bold = True
            title_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
            
            # Look for content under this heading
            i += 1
            content_lines = []
            
            # Collect all content until next heading
            while i < len(lines) and not lines[i].startswith('#'):
                current_line = lines[i].strip()
                
                if current_line.startswith('### '):  # Subsection
                    subsection_title = current_line[4:]
                    
                    # Add subsection slide
                    subsection_slide = prs.slides.add_slide(blank_slide_layout)
                    
                    # Add subsection title
                    sub_title_box = subsection_slide.shapes.add_textbox(Inches(1), Inches(1), Inches(11.33), Inches(1))
                    sub_title_tf = sub_title_box.text_frame
                    sub_title_tf.text = subsection_title
                    sub_title_tf.paragraphs[0].font.size = Pt(32)
                    sub_title_tf.paragraphs[0].font.bold = True
                    sub_title_tf.paragraphs[0].alignment = PP_ALIGN.LEFT
                    
                    # Collect items under subsection
                    i += 1
                    items = []
                    while i < len(lines) and not lines[i].startswith('#') and not lines[i].startswith('###'):
                        item_line = lines[i].strip()
                        if item_line.startswith('- '):
                            items.append(item_line[2:])
                        elif item_line.startswith('### ') or item_line.startswith('## ') or item_line.startswith('# '):
                            break
                        else:
                            # Add non-list content as regular paragraph
                            if item_line:
                                items.append(item_line)
                        i += 1
                    
                    # Add items to subsection slide
                    if items:
                        content_box = subsection_slide.shapes.add_textbox(Inches(1.5), Inches(2.5), Inches(10), Inches(4))
                        content_tf = content_box.text_frame
                        for idx, item in enumerate(items):
                            if idx == 0:
                                content_tf.paragraphs[0].text = item
                                content_tf.paragraphs[0].font.size = content_font_size
                            else:
                                p = content_tf.add_paragraph()
                                p.text = item
                                p.font.size = content_font_size
                            if item.startswith('- ') or item.startswith('* '):
                                p.level = 1
                        content_tf.word_wrap = True
                    
                    continue
                elif current_line.startswith('- '):  # List item
                    content_lines.append(current_line)
                elif current_line:  # Regular content
                    content_lines.append(current_line)
                
                i += 1
            
            # Add content slide if there's content
            if content_lines:
                content_slide = prs.slides.add_slide(blank_slide_layout)
                
                # Add title
                content_title_box = content_slide.shapes.add_textbox(Inches(1), Inches(1), Inches(11.33), Inches(0.8))
                content_title_tf = content_title_box.text_frame
                content_title_tf.text = section_title
                content_title_tf.paragraphs[0].font.size = heading_font_size
                content_title_tf.paragraphs[0].font.bold = True
                
                # Add content
                content_box = content_slide.shapes.add_textbox(Inches(1.5), Inches(2), Inches(10), Inches(5))
                content_tf = content_box.text_frame
                for idx, content_line in enumerate(content_lines):
                    if idx == 0:
                        content_tf.paragraphs[0].text = content_line
                        content_tf.paragraphs[0].font.size = content_font_size
                    else:
                        p = content_tf.add_paragraph()
                        p.text = content_line
                        p.font.size = content_font_size
                    if content_line.startswith('- ') or content_line.startswith('* '):
                        p.level = 1
                content_tf.word_wrap = True
        else:
            i += 1
    
    return prs

# Read markdown content from file
with open('/Users/xuan.lx/Documents/x-agent/x-agent/workspace/2026_ai_trends_report.md', 'r', encoding='utf-8') as f:
    md_content = f.read()

# Generate presentation
prs = create_presentation_from_markdown(md_content)

# Save presentation
output_path = '/Users/xuan.lx/Documents/x-agent/x-agent/workspace/presentations/2026_ai_trends.pptx'
prs.save(output_path)

print(f"Presentation saved to {output_path}")