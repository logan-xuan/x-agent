from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

def create_shanghai_travel_presentation():
    # Create presentation instance
    prs = Presentation()
    prs.slide_width = Inches(13.33)  # 16:9 aspect ratio
    prs.slide_height = Inches(7.5)

    # Define dark theme colors
    bg_color = RGBColor(44, 44, 44)      # Dark gray: #2C2C2C
    accent_color = RGBColor(74, 144, 226) # Blue: #4A90E2
    text_color = RGBColor(255, 255, 255)  # White: #FFFFFF
    sec_text_color = RGBColor(204, 204, 204)  # Light gray: #CCCCCC

    # Slide 1: Title Slide
    blank_slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(blank_slide_layout)
    
    # Background rectangle
    bg_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = bg_color
    bg_shape.line.color.rgb = bg_color
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.clear()
    p = title_frame.paragraphs[0]
    p.text = "上海魔都探索之旅"
    p.font.bold = True
    p.font.size = Pt(44)
    p.font.color.rgb = text_color
    p.alignment = PP_ALIGN.LEFT
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(11), Inches(0.8))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.clear()
    p = subtitle_frame.paragraphs[0]
    p.text = "杭州-上海七日游详细攻略"
    p.font.size = Pt(24)
    p.font.color.rgb = sec_text_color
    p.alignment = PP_ALIGN.LEFT
    
    # Footer
    footer_box = slide.shapes.add_textbox(Inches(1), Inches(7), Inches(11), Inches(0.5))
    footer_frame = footer_box.text_frame
    footer_frame.clear()
    p = footer_frame.paragraphs[0]
    p.text = "Travel Guide | 上海旅游"
    p.font.size = Pt(16)
    p.font.color.rgb = sec_text_color
    p.alignment = PP_ALIGN.CENTER

    # Slide 2: 行程概览
    slide = prs.slides.add_slide(blank_slide_layout)
    
    # Background
    bg_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = bg_color
    bg_shape.line.color.rgb = bg_color
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.clear()
    p = title_frame.paragraphs[0]
    p.text = "行程概览"
    p.font.bold = True
    p.font.size = Pt(36)
    p.font.color.rgb = text_color
    p.alignment = PP_ALIGN.LEFT
    
    # Content
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(6))
    content_frame = content_box.text_frame
    content_frame.clear()
    content_frame.word_wrap = True
    
    # Add days overview
    days = [
        "第一天: 抵达上海，游览外滩、南京路",
        "第二天: 上海迪士尼乐园一日游", 
        "第三天: 豫园及城隍庙，感受老上海风情",
        "第四天: 上海博物馆，陆家嘴金融区",
        "第五天: 新天地、田子坊文艺小资之旅",
        "第六天: 上海科技馆，世纪公园休闲游",
        "第七天: 自由活动，返程回家"
    ]
    
    for i, day in enumerate(days):
        if i > 0:
            p = content_frame.add_paragraph()
        else:
            p = content_frame.paragraphs[0]
        
        p.text = day
        p.font.size = Pt(18)
        p.font.color.rgb = sec_text_color
        p.space_after = Pt(12)
    
    # Footer
    footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(7), Inches(12), Inches(0.5))
    footer_frame = footer_box.text_frame
    footer_frame.clear()
    p = footer_frame.paragraphs[0]
    p.text = "Day 1-7 Overview"
    p.font.size = Pt(14)
    p.font.color.rgb = sec_text_color
    p.font.italic = True
    p.alignment = PP_ALIGN.LEFT

    # Slide 3: 第一天
    slide = prs.slides.add_slide(blank_slide_layout)
    
    # Background
    bg_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = bg_color
    bg_shape.line.color.rgb = bg_color
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.clear()
    p = title_frame.paragraphs[0]
    p.text = "第一天：经典上海印象"
    p.font.bold = True
    p.font.size = Pt(32)
    p.font.color.rgb = text_color
    p.alignment = PP_ALIGN.LEFT
    
    # Content with timeline
    timeline_items = [
        ("上午", "从杭州出发，乘坐高铁抵达上海（约1小时）"),
        ("下午", "游览外滩，欣赏万国建筑博览群"),
        ("傍晚", "漫步南京路步行街，体验繁华商业街"),
        ("晚上", "在外滩欣赏浦东陆家嘴夜景灯光秀")
    ]
    
    y_pos = 1.5
    for period, activity in timeline_items:
        # Time period
        time_box = slide.shapes.add_textbox(Inches(0.5), Inches(y_pos), Inches(2), Inches(0.5))
        time_frame = time_box.text_frame
        time_frame.clear()
        p = time_frame.paragraphs[0]
        p.text = period
        p.font.bold = True
        p.font.size = Pt(20)
        p.font.color.rgb = accent_color
        p.alignment = PP_ALIGN.LEFT
        
        # Activity
        activity_box = slide.shapes.add_textbox(Inches(2.7), Inches(y_pos), Inches(10), Inches(0.5))
        activity_frame = activity_box.text_frame
        activity_frame.clear()
        p = activity_frame.paragraphs[0]
        p.text = activity
        p.font.size = Pt(16)
        p.font.color.rgb = sec_text_color
        p.alignment = PP_ALIGN.LEFT
        
        y_pos += 0.7  # Increment Y position for next item

    # Slide 4: 第二天
    slide = prs.slides.add_slide(blank_slide_layout)
    
    # Background
    bg_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = bg_color
    bg_shape.line.color.rgb = bg_color
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.clear()
    p = title_frame.paragraphs[0]
    p.text = "第二天：梦幻迪士尼"
    p.font.bold = True
    p.font.size = Pt(32)
    p.font.color.rgb = text_color
    p.alignment = PP_ALIGN.LEFT
    
    # Content with timeline
    timeline_items = [
        ("全天", "上海迪士尼乐园一日游"),
        ("上午", "优先体验创极速光轮、雷鸣山漂流等热门项目"),
        ("下午", "观看花车巡游，体验各类主题园区"),
        ("晚上", "欣赏绚烂烟花表演结束完美一天")
    ]
    
    y_pos = 1.5
    for period, activity in timeline_items:
        # Time period
        time_box = slide.shapes.add_textbox(Inches(0.5), Inches(y_pos), Inches(2.5), Inches(0.5))
        time_frame = time_box.text_frame
        time_frame.clear()
        p = time_frame.paragraphs[0]
        p.text = period
        p.font.bold = True
        p.font.size = Pt(20)
        p.font.color.rgb = accent_color
        p.alignment = PP_ALIGN.LEFT
        
        # Activity
        activity_box = slide.shapes.add_textbox(Inches(2.7), Inches(y_pos), Inches(10), Inches(0.5))
        activity_frame = activity_box.text_frame
        activity_frame.clear()
        p = activity_frame.paragraphs[0]
        p.text = activity
        p.font.size = Pt(16)
        p.font.color.rgb = sec_text_color
        p.alignment = PP_ALIGN.LEFT
        
        y_pos += 0.7  # Increment Y position for next item

    # Slide 5: 第三天
    slide = prs.slides.add_slide(blank_slide_layout)
    
    # Background
    bg_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = bg_color
    bg_shape.line.color.rgb = bg_color
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.clear()
    p = title_frame.paragraphs[0]
    p.text = "第三天：老上海风情"
    p.font.bold = True
    p.font.size = Pt(32)
    p.font.color.rgb = text_color
    p.alignment = PP_ALIGN.LEFT
    
    # Content with timeline
    timeline_items = [
        ("上午", "游览豫园古典园林，欣赏江南园林之美"),
        ("中午", "在城隍庙小吃街品尝地道上海美食"),
        ("下午", "逛豫园商城，选购特色纪念品"),
        ("傍晚", "漫步十六铺码头，感受老上海码头文化")
    ]
    
    y_pos = 1.5
    for period, activity in timeline_items:
        # Time period
        time_box = slide.shapes.add_textbox(Inches(0.5), Inches(y_pos), Inches(2), Inches(0.5))
        time_frame = time_box.text_frame
        time_frame.clear()
        p = time_frame.paragraphs[0]
        p.text = period
        p.font.bold = True
        p.font.size = Pt(20)
        p.font.color.rgb = accent_color
        p.alignment = PP_ALIGN.LEFT
        
        # Activity
        activity_box = slide.shapes.add_textbox(Inches(2.7), Inches(y_pos), Inches(10), Inches(0.5))
        activity_frame = activity_box.text_frame
        activity_frame.clear()
        p = activity_frame.paragraphs[0]
        p.text = activity
        p.font.size = Pt(16)
        p.font.color.rgb = sec_text_color
        p.alignment = PP_ALIGN.LEFT
        
        y_pos += 0.7  # Increment Y position for next item

    # Slide 6: 第四天
    slide = prs.slides.add_slide(blank_slide_layout)
    
    # Background
    bg_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = bg_color
    bg_shape.line.color.rgb = bg_color
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.clear()
    p = title_frame.paragraphs[0]
    p.text = "第四天：文化与现代"
    p.font.bold = True
    p.font.size = Pt(32)
    p.font.color.rgb = text_color
    p.alignment = PP_ALIGN.LEFT
    
    # Content with timeline
    timeline_items = [
        ("上午", "参观上海博物馆，领略中华文明精华"),
        ("下午", "游览陆家嘴金融区，登上上海中心观景台"),
        ("傍晚", "在东方明珠塔或环球金融中心欣赏夜景"),
        ("晚上", "在陆家嘴享用高端江景晚餐")
    ]
    
    y_pos = 1.5
    for period, activity in timeline_items:
        # Time period
        time_box = slide.shapes.add_textbox(Inches(0.5), Inches(y_pos), Inches(2), Inches(0.5))
        time_frame = time_box.text_frame
        time_frame.clear()
        p = time_frame.paragraphs[0]
        p.text = period
        p.font.bold = True
        p.font.size = Pt(20)
        p.font.color.rgb = accent_color
        p.alignment = PP_ALIGN.LEFT
        
        # Activity
        activity_box = slide.shapes.add_textbox(Inches(2.7), Inches(y_pos), Inches(10), Inches(0.5))
        activity_frame = activity_box.text_frame
        activity_frame.clear()
        p = activity_frame.paragraphs[0]
        p.text = activity
        p.font.size = Pt(16)
        p.font.color.rgb = sec_text_color
        p.alignment = PP_ALIGN.LEFT
        
        y_pos += 0.7  # Increment Y position for next item

    # Slide 7: 第五天
    slide = prs.slides.add_slide(blank_slide_layout)
    
    # Background
    bg_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = bg_color
    bg_shape.line.color.rgb = bg_color
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.clear()
    p = title_frame.paragraphs[0]
    p.text = "第五天：文艺小资之旅"
    p.font.bold = True
    p.font.size = Pt(32)
    p.font.color.rgb = text_color
    p.alignment = PP_ALIGN.LEFT
    
    # Content with timeline
    timeline_items = [
        ("上午", "漫步新天地，感受石库门建筑改造典范"),
        ("中午", "在新天地品味精致本帮菜"),
        ("下午", "逛田子坊创意园区，探寻艺术小店"),
        ("晚上", "在田子坊享受创意料理或酒吧文化")
    ]
    
    y_pos = 1.5
    for period, activity in timeline_items:
        # Time period
        time_box = slide.shapes.add_textbox(Inches(0.5), Inches(y_pos), Inches(2), Inches(0.5))
        time_frame = time_box.text_frame
        time_frame.clear()
        p = time_frame.paragraphs[0]
        p.text = period
        p.font.bold = True
        p.font.size = Pt(20)
        p.font.color.rgb = accent_color
        p.alignment = PP_ALIGN.LEFT
        
        # Activity
        activity_box = slide.shapes.add_textbox(Inches(2.7), Inches(y_pos), Inches(10), Inches(0.5))
        activity_frame = activity_box.text_frame
        activity_frame.clear()
        p = activity_frame.paragraphs[0]
        p.text = activity
        p.font.size = Pt(16)
        p.font.color.rgb = sec_text_color
        p.alignment = PP_ALIGN.LEFT
        
        y_pos += 0.7  # Increment Y position for next item

    # Slide 8: 第六天
    slide = prs.slides.add_slide(blank_slide_layout)
    
    # Background
    bg_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = bg_color
    bg_shape.line.color.rgb = bg_color
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.clear()
    p = title_frame.paragraphs[0]
    p.text = "第六天：科普休闲之旅"
    p.font.bold = True
    p.font.size = Pt(32)
    p.font.color.rgb = text_color
    p.alignment = PP_ALIGN.LEFT
    
    # Content with timeline
    timeline_items = [
        ("上午", "参观上海科技馆，体验前沿科技魅力"),
        ("中午", "在科技馆附近用餐休息"),
        ("下午", "游览世纪公园，享受都市绿洲"),
        ("傍晚", "在世纪公园欣赏夕阳，或体验户外活动")
    ]
    
    y_pos = 1.5
    for period, activity in timeline_items:
        # Time period
        time_box = slide.shapes.add_textbox(Inches(0.5), Inches(y_pos), Inches(2), Inches(0.5))
        time_frame = time_box.text_frame
        time_frame.clear()
        p = time_frame.paragraphs[0]
        p.text = period
        p.font.bold = True
        p.font.size = Pt(20)
        p.font.color.rgb = accent_color
        p.alignment = PP_ALIGN.LEFT
        
        # Activity
        activity_box = slide.shapes.add_textbox(Inches(2.7), Inches(y_pos), Inches(10), Inches(0.5))
        activity_frame = activity_box.text_frame
        activity_frame.clear()
        p = activity_frame.paragraphs[0]
        p.text = activity
        p.font.size = Pt(16)
        p.font.color.rgb = sec_text_color
        p.alignment = PP_ALIGN.LEFT
        
        y_pos += 0.7  # Increment Y position for next item

    # Slide 9: 第七天和旅行贴士
    slide = prs.slides.add_slide(blank_slide_layout)
    
    # Background
    bg_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = bg_color
    bg_shape.line.color.rgb = bg_color
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.clear()
    p = title_frame.paragraphs[0]
    p.text = "第七天：返程 & 旅行贴士"
    p.font.bold = True
    p.font.size = Pt(32)
    p.font.color.rgb = text_color
    p.alignment = PP_ALIGN.LEFT
    
    # Timeline section
    timeline_items = [
        ("上午", "自由活动，购买特产"),
        ("下午", "前往火车站/机场，返回杭州")
    ]
    
    y_pos = 1.5
    for period, activity in timeline_items:
        # Time period
        time_box = slide.shapes.add_textbox(Inches(0.5), Inches(y_pos), Inches(2), Inches(0.5))
        time_frame = time_box.text_frame
        time_frame.clear()
        p = time_frame.paragraphs[0]
        p.text = period
        p.font.bold = True
        p.font.size = Pt(20)
        p.font.color.rgb = accent_color
        p.alignment = PP_ALIGN.LEFT
        
        # Activity
        activity_box = slide.shapes.add_textbox(Inches(2.7), Inches(y_pos), Inches(10), Inches(0.5))
        activity_frame = activity_box.text_frame
        activity_frame.clear()
        p = activity_frame.paragraphs[0]
        p.text = activity
        p.font.size = Pt(16)
        p.font.color.rgb = sec_text_color
        p.alignment = PP_ALIGN.LEFT
        
        y_pos += 0.7  # Increment Y position for next item
    
    # Travel tips title
    tips_title_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.2), Inches(12), Inches(0.6))
    tips_title_frame = tips_title_box.text_frame
    tips_title_frame.clear()
    p = tips_title_frame.paragraphs[0]
    p.text = "旅行贴士"
    p.font.bold = True
    p.font.size = Pt(20)
    p.font.color.rgb = accent_color
    p.alignment = PP_ALIGN.LEFT
    
    # Travel tips content
    tips_content_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(12), Inches(3))
    tips_content_frame = tips_content_box.text_frame
    tips_content_frame.clear()
    tips_content_frame.word_wrap = True
    
    tips = [
        "交通: 市内建议使用地铁，一日票50元无限次乘坐，便捷高效",
        "美食: 小笼包、生煎包、蟹粉小笼、白切鸡、糖醋排骨等上海本帮菜",
        "住宿: 建议选择人民广场、静安寺或陆家嘴附近，交通便利",
        "门票: 迪士尼平日票价435元，豫园门票40元，上海博物馆免费(需预约)"
    ]
    
    for i, tip in enumerate(tips):
        if i > 0:
            p = tips_content_frame.add_paragraph()
        else:
            p = tips_content_frame.paragraphs[0]
        
        p.text = f"• {tip}"
        p.font.size = Pt(16)
        p.font.color.rgb = sec_text_color
        p.space_after = Pt(8)

    # Slide 10: 上海景点概览
    slide = prs.slides.add_slide(blank_slide_layout)
    
    # Background
    bg_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = bg_color
    bg_shape.line.color.rgb = bg_color
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.clear()
    p = title_frame.paragraphs[0]
    p.text = "上海主要景点分布"
    p.font.bold = True
    p.font.size = Pt(32)
    p.font.color.rgb = text_color
    p.alignment = PP_ALIGN.LEFT
    
    # Content
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(6))
    content_frame = content_box.text_frame
    content_frame.clear()
    content_frame.word_wrap = True
    
    # Add location and attractions
    locations = [
        ("市中心区域:", [
            "• 外滩 - 万国建筑群，黄浦江夜景",
            "• 南京路步行街 - 百年商业街，繁华购物区",
            "• 人民广场 - 市中心地标，交通枢纽"
        ]),
        ("浦东新区:", [
            "• 陆家嘴 - 摩天大楼群，金融中心",
            "• 上海中心观景台 - 中国最高观景台",
            "• 东方明珠 - 上海地标建筑"
        ]),
        ("文化历史区:", [
            "• 豫园及城隍庙 - 江南古典园林，传统市井文化",
            "• 上海博物馆 - 中华文明艺术宝库",
            "• 新天地 - 石库门建筑群，时尚休闲区"
        ])
    ]
    
    for loc_name, attractions in locations:
        # Location name
        p = content_frame.add_paragraph() if content_frame.paragraphs != [content_frame.paragraphs[0]] else content_frame.paragraphs[0]
        p.text = loc_name
        p.font.bold = True
        p.font.size = Pt(16)
        p.font.color.rgb = accent_color
        p.space_after = Pt(4)
        
        # Attractions under this location
        for attraction in attractions:
            p = content_frame.add_paragraph()
            p.text = attraction
            p.font.size = Pt(16)
            p.font.color.rgb = sec_text_color
            p.space_after = Pt(4)

    # Save the presentation
    output_dir = 'presentations'
    os.makedirs(output_dir, exist_ok=True)
    output_path = os.path.join(output_dir, 'hangzhou_to_shanghai_travel_guide.pptx')
    prs.save(output_path)
    
    print(f"✅ PPT saved: {output_path}")
    return output_path

if __name__ == "__main__":
    create_shanghai_travel_presentation()