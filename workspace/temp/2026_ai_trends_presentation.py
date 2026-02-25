from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
import re

def create_2026_ai_trends_presentation():
    # Create presentation
    prs = Presentation()
    
    # Set slide dimensions (standard 16:9 aspect ratio)
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "2026年AI发展趋势深度研究报告"
    subtitle.text = "AI治理全球化 · 智能算力规模化 · 应用主流化\n多模态实用化 · 原生AI终端硬件普及化"
    
    # Customize title
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)  # Dark blue
    
    # Customize subtitle
    subtitle.text_frame.paragraphs[0].font.size = Pt(24)
    subtitle.text_frame.paragraphs[0].font.color.rgb = RGBColor(102, 102, 102)  # Gray
    
    # Add section slides and content based on markdown structure
    md_content = """# 2026年AI发展趋势深度研究报告

## 摘要
本报告深入分析了2026年AI发展的主要趋势，涵盖了AI治理全球化、智能算力规模化、应用主流化、多模态实用化、原生AI终端硬件普及化等多个方面，旨在为AI从业者提供全面的趋势洞察和战略指导。

## 1. 引言
2026年将是人工智能发展的重要节点，随着技术的不断成熟和应用场景的持续拓展，AI正从概念走向实践，从实验室走向产业一线。本报告将从多个维度深入分析2026年AI发展的核心趋势。

## 2. AI治理全球化趋势
### 2.1 国际合作加强
- 全球AI治理框架逐步完善
- 主要经济体间AI伦理标准趋同
- 跨国AI监管协调机制建立

### 2.2 法规政策演进
- 数据隐私保护法规升级
- AI算法透明度要求提高
- 自动化决策问责制度确立

## 3. 智能算力规模化发展
### 3.1 算力基础设施
- 专用AI芯片性能持续提升
- 云计算与边缘计算协同
- 算力资源共享平台兴起

### 3.2 算法效率优化
- 模型压缩技术成熟
- 知识蒸馏广泛应用
- 绿色AI理念普及

## 4. AI应用主流化
### 4.1 垂直行业渗透
- 医疗健康AI辅助诊断普及
- 金融风控智能化升级
- 制造业AI质检标准化

### 4.2 消费级应用爆发
- 个性化AI助手成为标配
- 智能家居生态深度融合
- AR/VR+AI体验革命

## 5. 多模态AI实用化
### 5.1 技术融合加速
- 视觉-语言模型能力增强
- 听觉-视觉交互体验优化
- 跨模态理解准确率提升

### 5.2 应用场景拓展
- 智能客服多模态交互
- 教育培训沉浸式体验
- 创意设计辅助工具

## 6. 原生AI终端硬件普及
### 6.1 设备智能化升级
- AI手机个人助理功能
- 智能眼镜AR显示技术
- 可穿戴设备健康监测

### 6.2 边缘AI芯片应用
- 低功耗AI处理器普及
- 端侧推理能力增强
- 隐私保护本地处理

## 7. 技术发展趋势预测
### 7.1 通用人工智能(GAGI)进展
- 认知架构统一性提升
- 推理能力跨领域迁移
- 自主学习机制完善

### 7.2 新兴技术融合
- 量子计算+AI突破
- 生物计算+AI探索
- 神经形态计算应用

## 8. 应用场景展望
### 8.1 智慧城市建设
- 交通流量智能调度
- 环境监测预警系统
- 公共安全风险评估

### 8.2 企业数字化转型
- 智能供应链管理
- 自动化业务流程
- 数据驱动决策支持

## 9. 挑战与机遇
### 9.1 主要挑战
- 数据质量与可用性问题
- 算法偏见与公平性考量
- 安全性与鲁棒性挑战

### 9.2 发展机遇
- 政策支持力度加大
- 投资市场持续关注
- 人才储备逐步充实

## 10. 结论与建议
2026年AI发展将迎来新高峰，建议企业提前布局关键技术，重视合规建设，培养复合型人才，抓住AI发展带来的历史机遇."""
    
    # Parse the markdown content and create slides
    sections = parse_markdown_to_sections(md_content)
    
    # Create content slides
    for i, section in enumerate(sections):
        if i == 0:  # Skip the first section as it's already covered in the title slide
            continue
            
        if section['level'] == 2:  # Main sections
            # Section title slide
            title_only_layout = prs.slide_layouts[1]  # Title and Content layout
            slide = prs.slides.add_slide(title_only_layout)
            
            title = slide.shapes.title
            title.text = section['title']
            title.text_frame.paragraphs[0].font.size = Pt(36)
            title.text_frame.paragraphs[0].font.bold = True
            title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)  # Dark blue
            
            # Add content if available
            if section['content']:
                left = Inches(1)
                top = Inches(2)
                width = Inches(11.33)
                height = Inches(4.5)
                
                textbox = slide.shapes.add_textbox(left, top, width, height)
                text_frame = textbox.text_frame
                text_frame.word_wrap = True
                
                # Process content into paragraphs
                paragraphs = section['content'].split('\n')
                first_para = True
                for para in paragraphs:
                    para = para.strip()
                    if para.startswith('###'):
                        # Subsection header
                        p = text_frame.add_paragraph()
                        p.text = para.replace('###', '').strip()
                        p.font.size = Pt(20)
                        p.font.bold = True
                        p.font.color.rgb = RGBColor(0, 102, 153)  # Medium blue
                        p.level = 1
                    elif para.startswith('- '):
                        # Bullet point
                        p = text_frame.add_paragraph()
                        p.text = para.replace('- ', '').strip()
                        p.font.size = Pt(16)
                        p.level = 2
                    elif para and not first_para:
                        # Regular paragraph
                        p = text_frame.add_paragraph()
                        p.text = para
                        p.font.size = Pt(16)
                        p.level = 1
                    elif first_para and para:
                        # First paragraph of section
                        text_frame.paragraphs[0].text = para
                        text_frame.paragraphs[0].font.size = Pt(16)
                        text_frame.paragraphs[0].level = 1
                        first_para = False
    
    # Final slide: Conclusion
    title_only_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(title_only_layout)
    
    title = slide.shapes.title
    title.text = "结论与建议"
    title.text_frame.paragraphs[0].font.size = Pt(36)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)  # Dark blue
    
    content = ("2026年AI发展将迎来新高峰，建议企业提前布局关键技术，\n"
               "重视合规建设，培养复合型人才，抓住AI发展带来的历史机遇。\n\n"
               "关键趋势总结：\n"
               "• AI治理全球化 - 国际合作与法规完善\n"
               "• 智能算力规模化 - 基础设施与效率优化\n"
               "• AI应用主流化 - 行业渗透与消费爆发\n"
               "• 多模态AI实用化 - 技术融合与场景拓展\n"
               "• 原生AI终端硬件普及 - 设备智能化与边缘计算")
    
    left = Inches(1)
    top = Inches(2)
    width = Inches(11.33)
    height = Inches(4.5)
    
    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.word_wrap = True
    text_frame.text = content
    
    for i, paragraph in enumerate(text_frame.paragraphs):
        if i == 0 or i == 1:  # Title paragraphs
            paragraph.font.size = Pt(18)
        elif i == 2:  # "关键趋势总结：" paragraph
            paragraph.font.size = Pt(20)
            paragraph.font.bold = True
        else:  # Bullet points
            paragraph.font.size = Pt(16)
            paragraph.level = 1
    
    return prs

def parse_markdown_to_sections(md_content):
    lines = md_content.split('\n')
    sections = []
    current_section = None
    
    for line in lines:
        line = line.strip()
        
        if line.startswith('## '):
            # New main section
            if current_section:
                sections.append(current_section)
            
            current_section = {
                'level': 2,
                'title': line[3:],  # Remove '## '
                'content': ''
            }
        elif line.startswith('### ') and current_section:
            # Subsection
            current_section['content'] += f"{line}\n"
        elif line and current_section and not line.startswith('#'):
            # Content for current section
            current_section['content'] += f"{line}\n"
    
    # Add the last section
    if current_section:
        sections.append(current_section)
    
    return sections

if __name__ == "__main__":
    prs = create_2026_ai_trends_presentation()
    prs.save('/Users/xuan.lx/Documents/x-agent/x-agent/workspace/presentations/2026_ai_trends.pptx')
    print("Presentation created successfully!")