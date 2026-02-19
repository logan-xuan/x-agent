from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

def create_chinese_new_year_travel_ppt():
    # Create presentation
    prs = Presentation()
    
    # Set slide width and height (16:9 aspect ratio)
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "春节旅游"
    subtitle.text = "欢度佳节，畅游天下\n祝您新春快乐，旅途愉快！"
    
    # Style title
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(220, 20, 60)  # Crimson red
    title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Style subtitle
    subtitle.text_frame.paragraphs[0].font.size = Pt(24)
    subtitle.text_frame.paragraphs[0].font.color.rgb = RGBColor(105, 105, 105)  # Dim gray
    subtitle.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Slide 2: Overview
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "春节旅游概述"
    content.text = ("春节是中国最重要的传统节日，也是家人团聚的时刻。\n\n" 
                   "随着生活水平的提高，越来越多的家庭选择在春节期间外出旅游，\n" 
                   "体验不同的年味和文化，创造美好的回忆。")

    # Slide 3: Domestic Destinations
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "热门国内旅游目的地"
    content.text = ("北方线路：\n"
                   "- 哈尔滨冰雪节：欣赏冰雕雪塑，感受北国风光\n"
                   "- 北京故宫：体验古都年味，参观皇家建筑\n\n"
                   "南方线路：\n"
                   "- 三亚：温暖海岛，避寒度假\n"
                   "- 西安：古都文化，历史之旅")

    # Slide 4: Overseas Travel
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "海外旅游推荐"
    content.text = ("亚洲：\n"
                   "- 日本：赏樱花、泡温泉、体验和式文化\n"
                   "- 泰国：热带风情、佛教文化、美食之旅\n\n"
                   "欧美：\n"
                   "- 新加坡：花园城市、多元文化\n"
                   "- 澳大利亚：反季节旅游、自然奇观")

    # Slide 5: Precautions
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "春节旅游注意事项"
    content.text = ("1. 提前预订：机票、酒店需提前预订，避免临时涨价\n"
                   "2. 天气准备：根据目的地天气情况准备衣物\n"
                   "3. 证件齐全：身份证、护照、签证等证件随身携带\n"
                   "4. 安全意识：注意人身和财产安全\n"
                   "5. 文化尊重：了解当地文化和习俗")

    # Slide 6: Budget Planning
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "预算规划与省钱技巧"
    content.text = ("预算分配建议：\n"
                   "- 交通费用：占总预算的30%\n"
                   "- 住宿费用：占总预算的25%\n"
                   "- 餐饮费用：占总预算的20%\n"
                   "- 景点门票：占总预算的15%\n"
                   "- 其他费用：占总预算的10%\n\n"
                   "省钱技巧：\n"
                   "- 错峰出行：避开高峰时段\n"
                   "- 团购优惠：利用团购平台获取折扣")

    # Slide 7: Special Experiences
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "春节特色体验活动"
    content.text = ("传统文化体验：\n"
                   "- 观看舞龙舞狮表演\n"
                   "- 逛庙会品尝特色小吃\n"
                   "- 参与包饺子、做年糕等活动\n\n"
                   "现代娱乐活动：\n"
                   "- 主题乐园过大年\n"
                   "- 温泉度假村放松身心\n"
                   "- 滑雪场享受冰雪乐趣")

    # Slide 8: Safety Tips
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "安全与健康提示"
    content.text = ("交通安全：\n"
                   "- 选择正规交通工具\n"
                   "- 遵守交通规则\n\n"
                   "饮食安全：\n"
                   "- 选择卫生条件良好的餐厅\n"
                   "- 注意食物新鲜程度\n\n"
                   "健康防护：\n"
                   "- 携带常用药品\n"
                   "- 注意保暖防寒")

    # Slide 9: Travel Checklist
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "旅行准备清单"
    content.text = ("证件类：\n"
                   "- 身份证、护照\n"
                   "- 签证（海外旅行）\n"
                   "- 各种预订确认单\n\n"
                   "生活用品：\n"
                   "- 衣物、洗漱用品\n"
                   "- 充电器、移动电源\n"
                   "- 常用药品")

    # Slide 10: Conclusion
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "结语与祝福"
    content.text = ("春节旅游是新时代的过年方式，让我们在旅途中感受不同地方的年味，\n"
                   "体验各地的风土人情，与家人朋友共同创造美好回忆。\n\n"
                   "祝愿大家春节旅途愉快，身体健康，阖家幸福！\n\n"
                   "新年快乐，万事如意！")

    # Save the presentation
    prs.save('/Users/xuan.lx/Documents/x-agent/x-agent/workspace/chinese_new_year_travel.pptx')
    print("Chinese New Year travel presentation created successfully!")

if __name__ == "__main__":
    create_chinese_new_year_travel_ppt()