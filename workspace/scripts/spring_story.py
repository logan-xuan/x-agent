from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

def create_spring_story_ppt():
    """Create a presentation about spring stories with appropriate design."""
    prs = Presentation()
    
    # Set up color scheme - Spring theme (light greens, soft pinks, sky blues)
    # Using SAGE & TERRACOTTA inspired palette from SKILL.md
    bg_color = RGBColor(0xF4, 0xF1, 0xDE)  # Cream background
    text_color = RGBColor(0x2C, 0x2C, 0x2C)  # Charcoal text
    accent_color = RGBColor(0x87, 0xA9, 0x6B)  # Sage green
    secondary_accent = RGBColor(0xE0, 0x7A, 0x5F)  # Terracotta orange
    
    # Slide 1: Title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "春天的故事"
    subtitle.text = "生机勃勃的季节"
    
    # Apply color scheme to title slide
    for shape in slide.shapes:
        if hasattr(shape, "text_frame"):
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(44) if shape == title else Pt(24)
                    run.font.color.rgb = text_color
                    if shape == title:
                        run.font.bold = True
    
    # Change background color
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = bg_color
    
    # Slide 2: Introduction
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "春天的序曲"
    content.text = "冬去春来的变化\n- 气温回升\n- 白昼变长\n- 生命复苏"
    
    # Apply formatting
    title.text_frame.paragraphs[0].font.size = Pt(36)
    title.text_frame.paragraphs[0].font.color.rgb = accent_color
    title.text_frame.paragraphs[0].font.bold = True
    
    content.text_frame.paragraphs[0].font.size = Pt(18)
    content.text_frame.paragraphs[0].font.color.rgb = text_color
    
    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = bg_color
    
    # Slide 3: Plants awakening
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "大地苏醒"
    content.text = "植物的复苏与生长\n- 芽苞初绽\n- 新绿满枝\n- 花蕾含苞待放"
    
    # Apply formatting
    title.text_frame.paragraphs[0].font.size = Pt(36)
    title.text_frame.paragraphs[0].font.color.rgb = accent_color
    title.text_frame.paragraphs[0].font.bold = True
    
    content.text_frame.paragraphs[0].font.size = Pt(18)
    content.text_frame.paragraphs[0].font.color.rgb = text_color
    
    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = bg_color
    
    # Slide 4: Animals returning
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "动物归来"
    content.text = "春天里的动物世界\n- 候鸟归巢\n- 昆虫苏醒\n- 动物繁殖季节"
    
    # Apply formatting
    title.text_frame.paragraphs[0].font.size = Pt(36)
    title.text_frame.paragraphs[0].font.color.rgb = accent_color
    title.text_frame.paragraphs[0].font.bold = True
    
    content.text_frame.paragraphs[0].font.size = Pt(18)
    content.text_frame.paragraphs[0].font.color.rgb = text_color
    
    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = bg_color
    
    # Slide 5: Spring rain
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "春雨滋润"
    content.text = "雨水对生命的意义\n- 滋润土壤\n- 催发新芽\n- 清洗尘埃"
    
    # Apply formatting
    title.text_frame.paragraphs[0].font.size = Pt(36)
    title.text_frame.paragraphs[0].font.color.rgb = accent_color
    title.text_frame.paragraphs[0].font.bold = True
    
    content.text_frame.paragraphs[0].font.size = Pt(18)
    content.text_frame.paragraphs[0].font.color.rgb = text_color
    
    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = bg_color
    
    # Slide 6: Blooming flowers
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "花开满园"
    content.text = "各种春花竞相绽放\n- 桃花\n- 樱花\n- 梨花\n- 迎春花"
    
    # Apply formatting
    title.text_frame.paragraphs[0].font.size = Pt(36)
    title.text_frame.paragraphs[0].font.color.rgb = accent_color
    title.text_frame.paragraphs[0].font.bold = True
    
    content.text_frame.paragraphs[0].font.size = Pt(18)
    content.text_frame.paragraphs[0].font.color.rgb = text_color
    
    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = bg_color
    
    # Slide 7: People celebrating spring
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "人们迎接春天"
    content.text = "春季活动与习俗\n- 踏青郊游\n- 种植花草\n- 放风筝\n- 春季清洁"
    
    # Apply formatting
    title.text_frame.paragraphs[0].font.size = Pt(36)
    title.text_frame.paragraphs[0].font.color.rgb = accent_color
    title.text_frame.paragraphs[0].font.bold = True
    
    content.text_frame.paragraphs[0].font.size = Pt(18)
    content.text_frame.paragraphs[0].font.color.rgb = text_color
    
    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = bg_color
    
    # Slide 8: Spring colors
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "春天的色彩"
    content.text = "自然界的缤纷色彩\n- 新绿\n- 淡粉\n- 天蓝\n- 柠檬黄"
    
    # Apply formatting
    title.text_frame.paragraphs[0].font.size = Pt(36)
    title.text_frame.paragraphs[0].font.color.rgb = accent_color
    title.text_frame.paragraphs[0].font.bold = True
    
    content.text_frame.paragraphs[0].font.size = Pt(18)
    content.text_frame.paragraphs[0].font.color.rgb = text_color
    
    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = bg_color
    
    # Slide 9: Symbolism of spring
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "春天的寓意"
    content.text = "新生与希望的象征\n- 新开始\n- 成长机会\n- 生命力\n- 乐观精神"
    
    # Apply formatting
    title.text_frame.paragraphs[0].font.size = Pt(36)
    title.text_frame.paragraphs[0].font.color.rgb = accent_color
    title.text_frame.paragraphs[0].font.bold = True
    
    content.text_frame.paragraphs[0].font.size = Pt(18)
    content.text_frame.paragraphs[0].font.color.rgb = text_color
    
    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = bg_color
    
    # Slide 10: Conclusion
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "感受春天的美好"
    content.text = "让我们一起\n欣赏春天的美景\n感受生命的活力\n拥抱新的开始"
    
    # Apply formatting
    title.text_frame.paragraphs[0].font.size = Pt(36)
    title.text_frame.paragraphs[0].font.color.rgb = accent_color
    title.text_frame.paragraphs[0].font.bold = True
    
    content.text_frame.paragraphs[0].font.size = Pt(20)
    content.text_frame.paragraphs[0].font.color.rgb = text_color
    content.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = bg_color
    
    # Create presentations directory if not exists
    output_dir = 'presentations'
    os.makedirs(output_dir, exist_ok=True)
    
    # Save the presentation
    output_path = os.path.join(output_dir, 'spring_story.pptx')
    prs.save(output_path)
    
    print(f"PPTX saved to: {output_path}")
    return output_path

if __name__ == "__main__":
    create_spring_story_ppt()