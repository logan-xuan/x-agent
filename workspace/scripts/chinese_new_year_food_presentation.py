from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

def create_chinese_new_year_food_presentation():
    # Create presentation
    prs = Presentation()
    
    # Set slide width and height (16:9 aspect ratio)
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "春节美食指南"
    subtitle.text = "传统年味，舌尖上的中国"
    
    # Customize title font
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(139, 0, 0)  # Dark red
    subtitle.text_frame.paragraphs[0].font.size = Pt(24)
    subtitle.text_frame.paragraphs[0].font.color.rgb = RGBColor(105, 105, 105)  # Gray
    
    # Slide 2: Introduction
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = '春节饮食文化'
    
    text_frame = body_shape.text_frame
    text_frame.clear()  # Clear default text
    p = text_frame.add_paragraph()
    p.text = "春节是中国最重要的传统节日，美食是其中不可或缺的一部分。"
    p.font_size = Pt(18)
    
    p = text_frame.add_paragraph()
    p.text = "每道菜都寓意着吉祥如意，寄托着人们对新年的美好祝愿。"
    p.font_size = Pt(16)
    p.level = 1
    
    # Slide 3: Main Dishes
    slide = prs.slides.add_slide(bullet_slide_layout)
    title_shape = slide.shapes.title
    body_shape = slide.placeholders[1]
    
    title_shape.text = '春节主食'
    
    text_frame = body_shape.text_frame
    text_frame.clear()
    p = text_frame.add_paragraph()
    p.text = "饺子 - 形似元宝，寓意招财进宝"
    p.font_size = Pt(18)
    
    p = text_frame.add_paragraph()
    p.text = "年糕 - 年年高升，生活更上一层楼"
    p.font_size = Pt(16)
    p.level = 1
    
    p = text_frame.add_paragraph()
    p.text = "汤圆 - 团团圆圆，家庭和睦"
    p.font_size = Pt(16)
    p.level = 1
    
    p = text_frame.add_paragraph()
    p.text = "春卷 - 迎春纳福，生机勃勃"
    p.font_size = Pt(16)
    p.level = 1
    
    # Slide 4: Fish Dishes
    slide = prs.slides.add_slide(bullet_slide_layout)
    title_shape = slide.shapes.title
    body_shape = slide.placeholders[1]
    
    title_shape.text = '鱼类菜肴'
    
    text_frame = body_shape.text_frame
    text_frame.clear()
    p = text_frame.add_paragraph()
    p.text = "年年有余 - 鱼是春节餐桌上的必备菜"
    p.font_size = Pt(18)
    
    p = text_frame.add_paragraph()
    p.text = "清蒸鲈鱼 - 肉质鲜嫩，老少皆宜"
    p.font_size = Pt(16)
    p.level = 1
    
    p = text_frame.add_paragraph()
    p.text = "红烧鲤鱼 - 寓意吉祥，色泽诱人"
    p.font_size = Pt(16)
    p.level = 1
    
    p = text_frame.add_paragraph()
    p.text = "糖醋鱼块 - 酸甜可口，开胃解腻"
    p.font_size = Pt(16)
    p.level = 1
    
    # Slide 5: Meat Dishes
    slide = prs.slides.add_slide(bullet_slide_layout)
    title_shape = slide.shapes.title
    body_shape = slide.placeholders[1]
    
    title_shape.text = '肉类佳肴'
    
    text_frame = body_shape.text_frame
    text_frame.clear()
    p = text_frame.add_paragraph()
    p.text = "春节餐桌上的丰盛肉类"
    p.font_size = Pt(18)
    
    p = text_frame.add_paragraph()
    p.text = "红烧肉 - 色泽红亮，寓意红红火火"
    p.font_size = Pt(16)
    p.level = 1
    
    p = text_frame.add_paragraph()
    p.text = "白切鸡 - 寓意吉利，肉质鲜美"
    p.font_size = Pt(16)
    p.level = 1
    
    p = text_frame.add_paragraph()
    p.text = "酱牛肉 - 营养丰富，寓意牛气冲天"
    p.font_size = Pt(16)
    p.level = 1
    
    # Slide 6: Vegetables and Healthy Options
    slide = prs.slides.add_slide(bullet_slide_layout)
    title_shape = slide.shapes.title
    body_shape = slide.placeholders[1]
    
    title_shape.text = '蔬菜搭配'
    
    text_frame = body_shape.text_frame
    text_frame.clear()
    p = text_frame.add_paragraph()
    p.text = "营养均衡的素食搭配"
    p.font_size = Pt(18)
    
    p = text_frame.add_paragraph()
    p.text = "炒青菜 - 清爽解腻，补充维生素"
    p.font_size = Pt(16)
    p.level = 1
    
    p = text_frame.add_paragraph()
    p.text = "凉拌菜 - 开胃小菜，口感清爽"
    p.font_size = Pt(16)
    p.level = 1
    
    p = text_frame.add_paragraph()
    p.text = "菌菇类 - 营养丰富，增强免疫力"
    p.font_size = Pt(16)
    p.level = 1
    
    # Slide 7: Desserts and Sweets
    slide = prs.slides.add_slide(bullet_slide_layout)
    title_shape = slide.shapes.title
    body_shape = slide.placeholders[1]
    
    title_shape.text = '甜品点心'
    
    text_frame = body_shape.text_frame
    text_frame.clear()
    p = text_frame.add_paragraph()
    p.text = "甜蜜幸福的新年甜品"
    p.font_size = Pt(18)
    
    p = text_frame.add_paragraph()
    p.text = "八宝饭 - 八宝齐聚，甜甜蜜蜜"
    p.font_size = Pt(16)
    p.level = 1
    
    p = text_frame.add_paragraph()
    p.text = "发糕 - 发财高升，寓意吉祥"
    p.font_size = Pt(16)
    p.level = 1
    
    p = text_frame.add_paragraph()
    p.text = "糖果瓜子 - 欢声笑语，甜蜜相伴"
    p.font_size = Pt(16)
    p.level = 1
    
    # Slide 8: Drinks
    slide = prs.slides.add_slide(bullet_slide_layout)
    title_shape = slide.shapes.title
    body_shape = slide.placeholders[1]
    
    title_shape.text = '节庆饮品'
    
    text_frame = body_shape.text_frame
    text_frame.clear()
    p = text_frame.add_paragraph()
    p.text = "春节餐桌上的特色饮品"
    p.font_size = Pt(18)
    
    p = text_frame.add_paragraph()
    p.text = "茶水 - 解腻助消化，清肠又暖胃"
    p.font_size = Pt(16)
    p.level = 1
    
    p = text_frame.add_paragraph()
    p.text = "黄酒/米酒 - 温润醇香，增添年味"
    p.font_size = Pt(16)
    p.level = 1
    
    p = text_frame.add_paragraph()
    p.text = "果汁饮料 - 老少皆宜，营养美味"
    p.font_size = Pt(16)
    p.level = 1
    
    # Slide 9: Regional Variations
    slide = prs.slides.add_slide(bullet_slide_layout)
    title_shape = slide.shapes.title
    body_shape = slide.placeholders[1]
    
    title_shape.text = '地域差异'
    
    text_frame = body_shape.text_frame
    text_frame.clear()
    p = text_frame.add_paragraph()
    p.text = "中国各地不同的春节饮食习俗"
    p.font_size = Pt(18)
    
    p = text_frame.add_paragraph()
    p.text = "北方 - 以饺子为主食"
    p.font_size = Pt(16)
    p.level = 1
    
    p = text_frame.add_paragraph()
    p.text = "南方 - 重视年糕和汤圆"
    p.font_size = Pt(16)
    p.level = 1
    
    p = text_frame.add_paragraph()
    p.text = "广东 - 白切鸡和盆菜"
    p.font_size = Pt(16)
    p.level = 1
    
    p = text_frame.add_paragraph()
    p.text = "四川 - 麻辣火锅和腊肉"
    p.font_size = Pt(16)
    p.level = 1
    
    # Slide 10: Conclusion
    conclusion_slide_layout = prs.slide_layouts[2]
    slide = prs.slides.add_slide(conclusion_slide_layout)
    title_shape = slide.shapes.title
    left_shape = slide.placeholders[1]
    right_shape = slide.placeholders[2]
    
    title_shape.text = '春节饮食小贴士'
    
    # Left text box
    left_shape.text = "健康饮食\n\n• 适量搭配荤素\n• 控制油盐糖摄入\n• 注意食物新鲜度"
    
    # Right text box
    right_shape.text = "安全用餐\n\n• 充分加热食物\n• 注意餐具清洁\n• 避免暴饮暴食"
    
    # Customize fonts
    for shape in [left_shape, right_shape]:
        for paragraph in shape.text_frame.paragraphs:
            paragraph.font.size = Pt(16)
    
    # Save the presentation
    file_path = "/Users/xuan.lx/Documents/x-agent/x-agent/workspace/chinese_new_year_food.pptx"
    prs.save(file_path)
    print(f"春节美食PPT已创建完成，保存在: {file_path}")

if __name__ == "__main__":
    create_chinese_new_year_food_presentation()