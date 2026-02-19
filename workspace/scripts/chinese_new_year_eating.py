from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

def create_chinese_new_year_eating_ppt():
    # Create presentation object with a blank slide layout
    prs = Presentation()
    
    # Add title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "春节美食文化"
    subtitle.text = "传统饮食习俗与现代庆祝方式"
    
    # Add content slides
    content_slides_data = [
        {
            "title": "春节饮食文化概述",
            "content": [
                "春节是中国最重要的传统节日",
                "饮食在春节期间具有特殊意义",
                "不同地区有不同的传统食物",
                "食物寓意吉祥和团圆"
            ]
        },
        {
            "title": "北方春节传统美食",
            "content": [
                "饺子 - 形似元宝，寓意招财进宝",
                "年糕 - 寓意年年高升",
                "春卷 - 象征春天的到来",
                "白切鸡 - 寓意吉祥如意"
            ]
        },
        {
            "title": "南方春节传统美食",
            "content": [
                "汤圆 - 寓意团团圆圆",
                "年糕 - 寓意年年高升",
                "发糕 - 寓意发财高升",
                "鱼 - 寓意年年有余"
            ]
        },
        {
            "title": "春节餐桌礼仪",
            "content": [
                "长辈先动筷",
                "不能用筷子敲碗",
                "鱼要完整上桌",
                "敬酒要双手捧杯",
                "说吉祥话助兴"
            ]
        },
        {
            "title": "现代春节聚餐趋势",
            "content": [
                "外出就餐越来越普遍",
                "融合菜系受到欢迎",
                "健康饮食观念增强",
                "素食主义者的选择增多"
            ]
        },
        {
            "title": "春节饮食注意事项",
            "content": [
                "适量进食，避免过量",
                "注意饮食均衡",
                "关注食品安全",
                "考虑特殊饮食需求"
            ]
        },
        {
            "title": "各地特色春节美食",
            "content": [
                "北京 - 糖葫芦、驴打滚",
                "上海 - 生煎包、小笼包",
                "广东 - 腊味煲仔饭、广式点心",
                "四川 - 腊肉香肠、火锅",
                "东北 - 杀猪菜、冻梨"
            ]
        },
        {
            "title": "春节饮品选择",
            "content": [
                "茶 - 传统饮品，有助消化",
                "红酒 - 现代聚会常见",
                "果汁 - 适合儿童饮用",
                "自制饮品 - 养生花茶"
            ]
        }
    ]
    
    # Add content slides
    for data in content_slides_data:
        bullet_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_slide_layout)
        
        title = slide.shapes.title
        title.text = data["title"]
        
        body_shape = slide.shapes.placeholders[1]
        tf = body_shape.text_frame
        
        for item in data["content"]:
            p = tf.add_paragraph()
            p.text = item
            p.level = 0
    
    # Add conclusion slide
    end_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(end_slide_layout)
    
    title = slide.shapes.title
    title.text = "春节饮食的美好寓意"
    
    body_shape = slide.shapes.placeholders[1]
    tf = body_shape.text_frame
    tf.text = "春节饮食不仅是一种生活方式，更是承载着深厚文化内涵的传统。每一道菜都寄托着人们对新年的美好祝愿。"
    
    p = tf.add_paragraph()
    p.text = "祝您春节快乐，阖家团圆！"
    p.font_size = Pt(18)
    p.alignment = PP_ALIGN.CENTER
    
    # Save the presentation
    filename = "chinese_new_year_eating.pptx"
    filepath = os.path.join(os.getcwd(), filename)
    prs.save(filepath)
    
    print(f"春节饮食文化PPT已保存至: {filepath}")
    return filepath

if __name__ == "__main__":
    create_chinese_new_year_eating_ppt()