from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.slide import SLIDE_LAYOUT_TYPE
import os

def create_travel_presentation():
    # Create presentation instance
    prs = Presentation()
    prs.slide_width = Inches(13.33)  # 16:9 aspect ratio
    prs.slide_height = Inches(7.5)

    # Define dark theme colors
    bg_color = RGBColor(44, 44, 44)      # Dark gray: #2C2C2C
    accent_color = RGBColor(74, 144, 226) # Blue: #4A90E2
    text_color = RGBColor(255, 255, 255)  # White: #FFFFFF
    sec_text_color = RGBColor(204, 204, 204)  # Light gray: #CCCCCC

    # Slide 1: Title Slide
    blank_slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(blank_slide_layout)
    
    # Background rectangle
    bg_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = bg_color
    bg_shape.line.color.rgb = bg_color
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.clear()
    p = title_frame.paragraphs[0]
    p.text = "西安古都探秘之旅"
    p.font.bold = True
    p.font.size = Pt(44)
    p.font.color.rgb = text_color
    p.alignment = PP_ALIGN.LEFT
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(11), Inches(0.8))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.clear()
    p = subtitle_frame.paragraphs[0]
    p.text = "杭州-西安五日游详细攻略"
    p.font.size = Pt(24)
    p.font.color.rgb = sec_text_color
    p.alignment = PP_ALIGN.LEFT
    
    # Footer
    footer_box = slide.shapes.add_textbox(Inches(1), Inches(7), Inches(11), Inches(0.5))
    footer_frame = footer_box.text_frame
    footer_frame.clear()
    p = footer_frame.paragraphs[0]
    p.text = "Travel Guide | 西安旅游"
    p.font.size = Pt(16)
    p.font.color.rgb = sec_text_color
    p.alignment = PP_ALIGN.CENTER

    # Slide 2: 行程概览
    slide = prs.slides.add_slide(blank_slide_layout)
    
    # Background
    bg_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = bg_color
    bg_shape.line.color.rgb = bg_color
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.clear()
    p = title_frame.paragraphs[0]
    p.text = "行程概览"
    p.font.bold = True
    p.font.size = Pt(36)
    p.font.color.rgb = text_color
    p.alignment = PP_ALIGN.LEFT
    
    # Content
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(6))
    content_frame = content_box.text_frame
    content_frame.clear()
    content_frame.word_wrap = True
    
    # Add days overview
    days = [
        "第一天: 抵达西安，市区游览 - 钟鼓楼、回民街",
        "第二天: 兵马俑 + 华清宫 - 世界文化遗产体验", 
        "第三天: 西安城墙 + 陕西历史博物馆 - 古都风貌",
        "第四天: 大雁塔 + 大兴善寺 - 盛唐文化",
        "第五天: 自由活动，返程回家"
    ]
    
    for i, day in enumerate(days):
        if i > 0:
            p = content_frame.add_paragraph()
        else:
            p = content_frame.paragraphs[0]
        
        p.text = day
        p.font.size = Pt(18)
        p.font.color.rgb = sec_text_color
        p.space_after = Pt(12)
    
    # Footer
    footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(7), Inches(12), Inches(0.5))
    footer_frame = footer_box.text_frame
    footer_frame.clear()
    p = footer_frame.paragraphs[0]
    p.text = "Day 1-5 Overview"
    p.font.size = Pt(14)
    p.font.color.rgb = sec_text_color
    p.font.italic = True
    p.alignment = PP_ALIGN.LEFT

    # Slide 3: 第一天
    slide = prs.slides.add_slide(blank_slide_layout)
    
    # Background
    bg_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = bg_color
    bg_shape.line.color.rgb = bg_color
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.clear()
    p = title_frame.paragraphs[0]
    p.text = "第一天：抵达西安"
    p.font.bold = True
    p.font.size = Pt(32)
    p.font.color.rgb = text_color
    p.alignment = PP_ALIGN.LEFT
    
    # Content with timeline
    timeline_items = [
        ("上午", "从杭州出发，乘机抵达西安"),
        ("下午", "入住酒店，稍作休息"),
        ("傍晚", "游览钟鼓楼广场，品尝回民街小吃")
    ]
    
    y_pos = 1.5
    for period, activity in timeline_items:
        # Time period
        time_box = slide.shapes.add_textbox(Inches(0.5), Inches(y_pos), Inches(2), Inches(0.5))
        time_frame = time_box.text_frame
        time_frame.clear()
        p = time_frame.paragraphs[0]
        p.text = period
        p.font.bold = True
        p.font.size = Pt(20)
        p.font.color.rgb = accent_color
        p.alignment = PP_ALIGN.LEFT
        
        # Activity
        activity_box = slide.shapes.add_textbox(Inches(2.7), Inches(y_pos), Inches(10), Inches(0.5))
        activity_frame = activity_box.text_frame
        activity_frame.clear()
        p = activity_frame.paragraphs[0]
        p.text = activity
        p.font.size = Pt(16)
        p.font.color.rgb = sec_text_color
        p.alignment = PP_ALIGN.LEFT
        
        y_pos += 0.7  # Increment Y position for next item
    
    # Slide 4: 第二天
    slide = prs.slides.add_slide(blank_slide_layout)
    
    # Background
    bg_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = bg_color
    bg_shape.line.color.rgb = bg_color
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.clear()
    p = title_frame.paragraphs[0]
    p.text = "第二天：世界文化遗产"
    p.font.bold = True
    p.font.size = Pt(32)
    p.font.color.rgb = text_color
    p.alignment = PP_ALIGN.LEFT
    
    # Content with timeline
    timeline_items = [
        ("上午", "参观秦始皇兵马俑博物馆 (游览约3小时)"),
        ("下午", "游览华清宫，感受唐玄宗与杨贵妃的爱情故事"),
        ("晚上 (可选)", "观看《长恨歌》实景演出 (298元起)")
    ]
    
    y_pos = 1.5
    for period, activity in timeline_items:
        # Time period
        time_box = slide.shapes.add_textbox(Inches(0.5), Inches(y_pos), Inches(2.5), Inches(0.5))
        time_frame = time_box.text_frame
        time_frame.clear()
        p = time_frame.paragraphs[0]
        p.text = period
        p.font.bold = True
        p.font.size = Pt(20)
        p.font.color.rgb = accent_color
        p.alignment = PP_ALIGN.LEFT
        
        # Activity
        activity_box = slide.shapes.add_textbox(Inches(2.7), Inches(y_pos), Inches(10), Inches(0.5))
        activity_frame = activity_box.text_frame
        activity_frame.clear()
        p = activity_frame.paragraphs[0]
        p.text = activity
        p.font.size = Pt(16)
        p.font.color.rgb = sec_text_color
        p.alignment = PP_ALIGN.LEFT
        
        y_pos += 0.7  # Increment Y position for next item

    # Slide 5: 第三天
    slide = prs.slides.add_slide(blank_slide_layout)
    
    # Background
    bg_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = bg_color
    bg_shape.line.color.rgb = bg_color
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.clear()
    p = title_frame.paragraphs[0]
    p.text = "第三天：古都风貌"
    p.font.bold = True
    p.font.size = Pt(32)
    p.font.color.rgb = text_color
    p.alignment = PP_ALIGN.LEFT
    
    # Content with timeline
    timeline_items = [
        ("上午", "参观西安城墙，骑行或步行体验明代城垣"),
        ("下午", "游览陕西历史博物馆 (免费，需预约)"),
        ("晚上", "游览大唐不夜城，感受盛唐文化氛围")
    ]
    
    y_pos = 1.5
    for period, activity in timeline_items:
        # Time period
        time_box = slide.shapes.add_textbox(Inches(0.5), Inches(y_pos), Inches(2), Inches(0.5))
        time_frame = time_box.text_frame
        time_frame.clear()
        p = time_frame.paragraphs[0]
        p.text = period
        p.font.bold = True
        p.font.size = Pt(20)
        p.font.color.rgb = accent_color
        p.alignment = PP_ALIGN.LEFT
        
        # Activity
        activity_box = slide.shapes.add_textbox(Inches(2.7), Inches(y_pos), Inches(10), Inches(0.5))
        activity_frame = activity_box.text_frame
        activity_frame.clear()
        p = activity_frame.paragraphs[0]
        p.text = activity
        p.font.size = Pt(16)
        p.font.color.rgb = sec_text_color
        p.alignment = PP_ALIGN.LEFT
        
        y_pos += 0.7  # Increment Y position for next item

    # Slide 6: 第四天
    slide = prs.slides.add_slide(blank_slide_layout)
    
    # Background
    bg_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = bg_color
    bg_shape.line.color.rgb = bg_color
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.clear()
    p = title_frame.paragraphs[0]
    p.text = "第四天：盛唐文化"
    p.font.bold = True
    p.font.size = Pt(32)
    p.font.color.rgb = text_color
    p.alignment = PP_ALIGN.LEFT
    
    # Content with timeline
    timeline_items = [
        ("上午", "参观大雁塔及北广场音乐喷泉"),
        ("下午", "游览大兴善寺，体验佛教文化"),
        ("傍晚", "参观大唐芙蓉园，欣赏园林景观")
    ]
    
    y_pos = 1.5
    for period, activity in timeline_items:
        # Time period
        time_box = slide.shapes.add_textbox(Inches(0.5), Inches(y_pos), Inches(2), Inches(0.5))
        time_frame = time_box.text_frame
        time_frame.clear()
        p = time_frame.paragraphs[0]
        p.text = period
        p.font.bold = True
        p.font.size = Pt(20)
        p.font.color.rgb = accent_color
        p.alignment = PP_ALIGN.LEFT
        
        # Activity
        activity_box = slide.shapes.add_textbox(Inches(2.7), Inches(y_pos), Inches(10), Inches(0.5))
        activity_frame = activity_box.text_frame
        activity_frame.clear()
        p = activity_frame.paragraphs[0]
        p.text = activity
        p.font.size = Pt(16)
        p.font.color.rgb = sec_text_color
        p.alignment = PP_ALIGN.LEFT
        
        y_pos += 0.7  # Increment Y position for next item

    # Slide 7: 第五天和旅行贴士
    slide = prs.slides.add_slide(blank_slide_layout)
    
    # Background
    bg_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = bg_color
    bg_shape.line.color.rgb = bg_color
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.clear()
    p = title_frame.paragraphs[0]
    p.text = "第五天：返程 & 旅行贴士"
    p.font.bold = True
    p.font.size = Pt(32)
    p.font.color.rgb = text_color
    p.alignment = PP_ALIGN.LEFT
    
    # Timeline section
    timeline_items = [
        ("上午", "自由活动，购买特产"),
        ("下午", "前往机场，返回杭州")
    ]
    
    y_pos = 1.5
    for period, activity in timeline_items:
        # Time period
        time_box = slide.shapes.add_textbox(Inches(0.5), Inches(y_pos), Inches(2), Inches(0.5))
        time_frame = time_box.text_frame
        time_frame.clear()
        p = time_frame.paragraphs[0]
        p.text = period
        p.font.bold = True
        p.font.size = Pt(20)
        p.font.color.rgb = accent_color
        p.alignment = PP_ALIGN.LEFT
        
        # Activity
        activity_box = slide.shapes.add_textbox(Inches(2.7), Inches(y_pos), Inches(10), Inches(0.5))
        activity_frame = activity_box.text_frame
        activity_frame.clear()
        p = activity_frame.paragraphs[0]
        p.text = activity
        p.font.size = Pt(16)
        p.font.color.rgb = sec_text_color
        p.alignment = PP_ALIGN.LEFT
        
        y_pos += 0.7  # Increment Y position for next item
    
    # Travel tips title
    tips_title_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.2), Inches(12), Inches(0.6))
    tips_title_frame = tips_title_box.text_frame
    tips_title_frame.clear()
    p = tips_title_frame.paragraphs[0]
    p.text = "旅行贴士"
    p.font.bold = True
    p.font.size = Pt(20)
    p.font.color.rgb = accent_color
    p.alignment = PP_ALIGN.LEFT
    
    # Travel tips content
    tips_content_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(12), Inches(3))
    tips_content_frame = tips_content_box.text_frame
    tips_content_frame.clear()
    tips_content_frame.word_wrap = True
    
    tips = [
        "交通: 市内建议使用地铁和出租车，去临潼可乘坐旅游专线",
        "美食: 肉夹馍、凉皮、羊肉泡馍、biáng biáng面、葫芦头等",
        "住宿: 建议选择钟楼、大雁塔附近，交通便利",
        "门票: 兵马俑120元，华清宫120元，西安城墙54元，陕历博免费(需预约)"
    ]
    
    for i, tip in enumerate(tips):
        if i > 0:
            p = tips_content_frame.add_paragraph()
        else:
            p = tips_content_frame.paragraphs[0]
        
        p.text = f"• {tip}"
        p.font.size = Pt(16)
        p.font.color.rgb = sec_text_color
        p.space_after = Pt(8)

    # Slide 8: 西安景点概览
    slide = prs.slides.add_slide(blank_slide_layout)
    
    # Background
    bg_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = bg_color
    bg_shape.line.color.rgb = bg_color
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.clear()
    p = title_frame.paragraphs[0]
    p.text = "西安主要景点分布"
    p.font.bold = True
    p.font.size = Pt(32)
    p.font.color.rgb = text_color
    p.alignment = PP_ALIGN.LEFT
    
    # Content
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(6))
    content_frame = content_box.text_frame
    content_frame.clear()
    content_frame.word_wrap = True
    
    # Add location and attractions
    locations = [
        ("市中心区域:", [
            "• 钟鼓楼广场 - 古都地标，夜景迷人",
            "• 回民街 - 美食天堂，体验当地小吃",
            "• 西安城墙 - 保存最完整的古代城垣"
        ]),
        ("东部临潼区:", [
            "• 兵马俑 - 世界第八大奇迹",
            "• 华清宫 - 唐玄宗与杨贵妃的爱情地"
        ]),
        ("南部文化区:", [
            "• 陕西历史博物馆 - 陕西历史文化宝库",
            "• 大雁塔 - 唐代著名佛塔",
            "• 大唐不夜城 - 盛唐文化体验区"
        ])
    ]
    
    for loc_name, attractions in locations:
        # Location name
        p = content_frame.add_paragraph() if content_frame.paragraphs != [content_frame.paragraphs[0]] else content_frame.paragraphs[0]
        p.text = loc_name
        p.font.bold = True
        p.font.size = Pt(16)
        p.font.color.rgb = accent_color
        p.space_after = Pt(4)
        
        # Attractions under this location
        for attraction in attractions:
            p = content_frame.add_paragraph()
            p.text = attraction
            p.font.size = Pt(16)
            p.font.color.rgb = sec_text_color
            p.space_after = Pt(4)

    # Save the presentation
    output_dir = 'presentations'
    os.makedirs(output_dir, exist_ok=True)
    output_path = os.path.join(output_dir, 'hangzhou_to_xian_travel_guide.pptx')
    prs.save(output_path)
    
    print(f"✅ PPT saved: {output_path}")
    return output_path

if __name__ == "__main__":
    create_travel_presentation()