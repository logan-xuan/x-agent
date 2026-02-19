from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

def create_braised_pork_belly_ppt():
    """Create a presentation about how to make braised pork belly (红烧肉)."""
    prs = Presentation()
    
    # Set slide dimensions to 16:9
    prs.slide_width = Inches(13.33)  # 720pt
    prs.slide_height = Inches(7.5)   # 405pt
    
    # Title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "红烧肉制作指南"
    subtitle.text = "经典中式美食 - 步骤详解"
    
    # Slide 2: Ingredients
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "所需食材"
    ingredients = [
        "五花肉 500克",
        "冰糖 30克",
        "生抽 2汤匙",
        "老抽 1汤匙",
        "料酒 2汤匙",
        "生姜 3片",
        "葱 2根",
        "八角 2颗",
        "桂皮 1小段"
    ]
    content.text = "\n".join(ingredients)
    
    # Slide 3: Preparation
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "准备工作"
    prep_steps = [
        "1. 五花肉切成3厘米见方的块",
        "2. 冷水下锅焯水，去除血沫",
        "3. 捞出洗净备用",
        "4. 准备调料：姜切片，葱打结"
    ]
    content.text = "\n".join(prep_steps)
    
    # Slide 4: Cooking Steps (Part 1)
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "烹饪步骤 - 第一部分"
    cooking_steps_1 = [
        "1. 热锅冷油，放入冰糖炒制糖色",
        "2. 糖色呈焦糖色时下入五花肉",
        "3. 大火翻炒至肉块均匀上色",
        "4. 加入料酒继续翻炒"
    ]
    content.text = "\n".join(cooking_steps_1)
    
    # Slide 5: Cooking Steps (Part 2)
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "烹饪步骤 - 第二部分"
    cooking_steps_2 = [
        "5. 加入生抽、老抽调色调味",
        "6. 加入足量热水没过肉块",
        "7. 放入葱姜和香料",
        "8. 大火烧开后转小火慢炖"
    ]
    content.text = "\n".join(cooking_steps_2)
    
    # Slide 6: Finishing Steps
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "收汁完成"
    finishing_steps = [
        "9. 慢炖约1小时至肉质软烂",
        "10. 开盖转大火收汁",
        "11. 汤汁浓稠即可关火",
        "12. 装盘撒上葱花装饰"
    ]
    content.text = "\n".join(finishing_steps)
    
    # Slide 7: Tips
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "制作要点"
    tips = [
        "• 选用肥瘦相间的五花肉",
        "• 炒糖色时火候要控制好",
        "• 慢炖让肉质更加软糯",
        "• 最后的收汁步骤很重要"
    ]
    content.text = "\n".join(tips)
    
    # Slide 8: Serving Suggestion
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "搭配建议"
    serving_suggestions = [
        "• 搭配米饭享用",
        "• 配以青菜解腻",
        "• 可加入土豆或鹌鹑蛋同煮",
        "• 是节庆餐桌上的经典菜肴"
    ]
    content.text = "\n".join(serving_suggestions)
    
    # Slide 9: History & Culture
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "文化背景"
    history_info = [
        "• 红烧肉是中国经典家常菜",
        "• 起源于江南地区",
        "• 毛泽东主席喜爱的菜品",
        "• 体现了中华饮食文化的精髓"
    ]
    content.text = "\n".join(history_info)
    
    # Slide 10: Thank You
    slide = prs.slides.add_slide(prs.slide_layouts[2])
    title = slide.shapes.title
    title.text = "谢谢观看"
    
    subtitle = slide.placeholders[1]
    subtitle.text = "祝您做出美味的红烧肉！\n\n欢迎品尝中国传统美食"
    
    # CRITICAL: Save to presentations/ directory, NOT workspace root!
    output_dir = 'presentations'
    os.makedirs(output_dir, exist_ok=True)  # Create directory if needed
    
    output_path = os.path.join(output_dir, 'braised_pork_belly_recipe.pptx')
    prs.save(output_path)
    
    print(f"✅ PPTX saved to: {output_path}")
    return output_path

if __name__ == "__main__":
    create_braised_pork_belly_ppt()
