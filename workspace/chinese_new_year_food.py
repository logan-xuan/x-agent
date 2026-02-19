from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

def create_chinese_new_year_food_ppt():
    # Create presentation
    prs = Presentation()
    
    # Set slide width and height (16:9 aspect ratio)
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # Title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "春节美食文化"
    subtitle.text = "传统年味，舌尖上的中国年"
    
    # Slide 2: Spring Festival Food Overview
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "春节饮食文化概览"
    content.text = """春节作为中国传统最重要的节日，饮食文化承载着深厚的历史底蕴和美好寓意。
    
    • 食物不仅是味蕾享受，更是文化传承
    • 每道菜都有吉祥寓意
    • 家庭团圆饭是核心环节"""
    
    # Slide 3: Northern Traditional Foods
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "北方传统美食"
    content.text = """• 饺子 - 形似元宝，寓意招财进宝
    • 饺子馅料丰富多样：韭菜鸡蛋、猪肉大葱、三鲜等
    • 年夜饭必吃，代表团圆和美满
    
    • 饽饽 - 各式面食，象征富足
    • 猪肉炖粉条 - 丰盛美味"""
    
    # Slide 4: Southern Traditional Foods
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "南方传统美食"
    content.text = """• 年糕 - 寓意年年高升
    • 米制糕点，口感软糯香甜
    • 南方春节必备食品
    
    • 汤圆 - 团团圆圆，甜甜蜜蜜
    • 寓意家庭和睦，生活圆满
    • 元宵节主要食品"""
    
    # Slide 5: Cantonese New Year Dishes
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "广府新年菜肴"
    content.text = """• 白切鸡 - 寓意有计(吉)有余
    • 整只烹制，完整上桌，象征完整圆满
    
    • 腊味合蒸 - 寓意丰收富足
    • 包含腊肠、腊肉、腊鸭等
    • 香味浓郁，风味独特
    
    • 发菜蚝豉 - 寓意发财好事"""
    
    # Slide 6: Symbolic Meanings
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "食物的吉祥寓意"
    content.text = """• 鱼 - 年年有余
• 饺子 - 招财进宝
• 年糕 - 年年高升
• 汤圆 - 团团圆圆
• 饺子形似元宝 - 财源滚滚
• 长寿面 - 长命百岁
• 猪蹄 - 待遇从厚"""
    
    # Slide 7: Regional Variations
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "地域差异特色"
    content.text = """• 北方：以面食为主，饺子是主角
• 南方：米制品占主导，年糕汤圆常见
• 江浙：精致小食，甜品丰富
• 川渝：麻辣口味，火锅盛行
• 广东：海鲜丰富，腊味多样
• 东北：炖菜丰盛，分量十足"""
    
    # Slide 8: Modern Celebrations
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "现代春节饮食"
    content.text = """• 传统与创新结合
• 外卖年夜饭兴起
• 旅游过年中的地方美食体验
• 健康饮食理念融入传统
• 网购年货便捷化
• 素食主义新趋势"""
    
    # Slide 9: Table Manners and Customs
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "餐桌礼仪与习俗"
    content.text = """• 长辈先动筷
• 不可翻拣菜肴
• 鱼不能翻面，寓意翻船
• 饺子剩几个，寓意年年有余
• 敬酒顺序有讲究
• 筷子不可插在米饭上
• 不说不吉利的话"""
    
    # Slide 10: Health Considerations
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "健康饮食建议"
    content.text = """• 荤素搭配，营养均衡
• 控制油炸食品摄入
• 多喝温水，助消化
• 适量饮酒，避免过量
• 注意食品安全
• 合理安排用餐时间
• 适当运动，促进消化"""
    
    # Slide 11: Conclusion
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "总结"
    content.text = """春节饮食不仅是味觉享受，更是文化传承的重要载体。
    
    • 传统美食承载美好寓意
    • 地域特色丰富多彩
    • 家庭团圆的重要纽带
    • 文化传承的重要方式
    
    祝大家春节快乐，龙年大吉！"""
    
    # Save the presentation
    filename = "chinese_new_year_food_ppt.pptx"
    filepath = os.path.join(os.getcwd(), filename)
    prs.save(filepath)
    print(f"PPT created successfully at {filepath}")
    
    return filepath

if __name__ == "__main__":
    create_chinese_new_year_food_ppt()