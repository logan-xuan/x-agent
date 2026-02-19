from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import os

def create_ai_agent_presentation():
    """Create a presentation about designing AI agent applications"""
    
    # Create a new presentation
    prs = Presentation()
    
    # Slide 1: Title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Designing AI Agent Applications"
    subtitle.text = "A Comprehensive Guide to Building Intelligent Systems\nCreated by 虾铁蛋"
    
    # Slide 2: What are AI Agents?
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = 'What are AI Agents?'
    
    tf = body_shape.text_frame
    tf.text = 'Definition:'
    
    p = tf.add_paragraph()
    p.text = 'AI Agents are autonomous systems that perceive their environment and take actions to achieve specific goals'
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = 'Key characteristics: Autonomy, Reactivity, Proactivity, Social Ability'
    p.level = 1
    
    # Slide 3: Types of AI Agents
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = 'Types of AI Agents'
    
    tf = body_shape.text_frame
    tf.text = 'Simple Reflex Agents'
    
    p = tf.add_paragraph()
    p.text = 'Model-based Reflex Agents'
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = 'Goal-based Agents'
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = 'Utility-based Agents'
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = 'Learning Agents'
    p.level = 1
    
    # Slide 4: Design Principles
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = 'Key Design Principles'
    
    tf = body_shape.text_frame
    tf.text = 'Modularity'
    
    p = tf.add_paragraph()
    p.text = 'Separate perception, decision-making, and action components'
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = 'Scalability'
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = 'Ensure the system can handle increasing complexity'
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = 'Robustness'
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = 'Handle unexpected situations gracefully'
    p.level = 1
    
    # Slide 5: Architecture Components
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = 'Core Architecture Components'
    
    tf = body_shape.text_frame
    tf.text = 'Perception Module'
    
    p = tf.add_paragraph()
    p.text = 'Processes input data from environment sensors'
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = 'Decision Engine'
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = 'Processes state and determines optimal actions'
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = 'Action Executor'
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = 'Performs actions in the environment'
    p.level = 1
    
    # Slide 6: Implementation Strategies
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = 'Implementation Strategies'
    
    tf = body_shape.text_frame
    tf.text = 'Rule-based Systems'
    
    p = tf.add_paragraph()
    p.text = 'Define explicit rules for agent behavior'
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = 'Machine Learning Models'
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = 'Train agents using supervised, unsupervised, or reinforcement learning'
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = 'Hybrid Approaches'
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = 'Combine multiple techniques for complex behaviors'
    p.level = 1
    
    # Slide 7: Best Practices
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = 'Best Practices for AI Agent Development'
    
    tf = body_shape.text_frame
    tf.text = 'Start Simple'
    
    p = tf.add_paragraph()
    p.text = 'Begin with basic functionality and gradually add complexity'
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = 'Test Thoroughly'
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = 'Validate agent behavior under various conditions'
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = 'Monitor Performance'
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = 'Continuously track and improve agent effectiveness'
    p.level = 1
    
    # Slide 8: Future Trends
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = 'Future Trends in AI Agent Development'
    
    tf = body_shape.text_frame
    tf.text = 'Multi-Agent Systems'
    
    p = tf.add_paragraph()
    p.text = 'Coordination between multiple intelligent agents'
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = 'Explainable AI'
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = 'Making agent decisions transparent and interpretable'
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = 'Human-AI Collaboration'
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = 'Seamless integration of human and artificial intelligence'
    p.level = 1
    
    # Slide 9: Conclusion
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = 'Conclusion'
    
    tf = body_shape.text_frame
    tf.text = 'AI agents represent the future of intelligent automation'
    
    p = tf.add_paragraph()
    p.text = 'Successful design requires understanding both technical and practical aspects'
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = 'Focus on modularity, scalability, and robustness'
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = 'Stay updated with emerging trends and technologies'
    p.level = 1
    
    # Save the presentation
    filename = 'ai_agent_design_presentation.pptx'
    filepath = os.path.join(os.getcwd(), filename)
    prs.save(filepath)
    print(f"Presentation created successfully: {filepath}")
    return filepath

if __name__ == "__main__":
    create_ai_agent_presentation()