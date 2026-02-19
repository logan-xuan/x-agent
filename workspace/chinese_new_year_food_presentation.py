from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

def create_chinese_new_year_food_presentation():
    # Create presentation
    prs = Presentation()
    
    # Set slide width and height (standard 4:3 aspect ratio)
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # Title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "春节美食文化"
    subtitle.text = "传统年味，舌尖上的团圆"
    
    # Customize title font
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.bold = True
    
    # Slide 1: 春节饮食文化概述
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = '春节饮食文化概述'
    
    tf = body_shape.text_frame
    tf.text = '春节饮食的意义'
    p = tf.add_paragraph()
    p.text = '• 团圆饭 - 家庭团聚的象征'
    p.level = 1
    p = tf.add_paragraph()
    p.text = '• 年菜寓意 - 吉祥如意的美好祝愿'
    p.level = 1
    p = tf.add_paragraph()
    p.text = '• 传统习俗 - 代代相传的文化符号'
    p.level = 1
    
    # Slide 2: 北方春节美食
    slide = prs.slides.add_slide(bullet_slide_layout)
    title_shape = slide.shapes.title
    body_shape = slide.placeholders[1]
    
    title_shape.text = '北方春节美食'
    
    tf = body_shape.text_frame
    tf.text = '饺子'
    p = tf.add_paragraph()
    p.text = '• 形状像元宝，寓意招财进宝'
    p.level = 1
    p = tf.add_paragraph()
    p.text = '• 年三十包饺子，大年初一吃饺子'
    p.level = 1
    p = tf.add_paragraph()
    p.text = '年糕'
    p.level = 1
    p = tf.add_paragraph()
    p.text = '• 寓意年年高升'
    p.level = 2
    p = tf.add_paragraph()
    p.text = '面条'
    p.level = 1
    p = tf.add_paragraph()
    p.text = '• 长寿面，寓意长命百岁'
    p.level = 2
    
    # Slide 3: 南方春节美食
    slide = prs.slides.add_slide(bullet_slide_layout)
    title_shape = slide.shapes.title
    body_shape = slide.placeholders[1]
    
    title_shape.text = '南方春节美食'
    
    tf = body_shape.text_frame
    tf.text = '汤圆'
    p = tf.add_paragraph()
    p.text = '• 圆形代表团圆美满'
    p.level = 1
    p = tf.add_paragraph()
    p.text = '• 元宵节必吃食品'
    p.level = 1
    p = tf.add_paragraph()
    p.text = '年糕'
    p.level = 1
    p = tf.add_paragraph()
    p.text = '• 糯米制作，软糯香甜'
    p.level = 2
    p = tf.add_paragraph()
    p.text = '春卷'
    p.level = 1
    p = tf.add_paragraph()
    p.text = '• 外酥内嫩，寓意迎春纳福'
    p.level = 2
    
    # Slide 4: 春节餐桌必备菜品
    slide = prs.slides.add_slide(bullet_slide_layout)
    title_shape = slide.shapes.title
    body_shape = slide.placeholders[1]
    
    title_shape.text = '春节餐桌必备菜品'
    
    tf = body_shape.text_frame
    tf.text = '鱼'
    p = tf.add_paragraph()
    p.text = '• 年年有余的美好寓意'
    p.level = 1
    p = tf.add_paragraph()
    p.text = '• 红烧鱼、清蒸鱼等多种做法'
    p.level = 1
    p = tf.add_paragraph()
    p.text = '鸡'
    p.level = 1
    p = tf.add_paragraph()
    p.text = '• 寓意吉利，大吉大利'
    p.level = 2
    p = tf.add_paragraph()
    p.text = '肘子'
    p.level = 1
    p = tf.add_paragraph()
    p.text = '• 肉食代表，丰盛富足'
    p.level = 2
    
    # Slide 5: 春节特色小吃
    slide = prs.slides.add_slide(bullet_slide_layout)
    title_shape = slide.shapes.title
    body_shape = slide.placeholders[1]
    
    title_shape.text = '春节特色小吃'
    
    tf = body_shape.text_frame
    tf.text = '糖果'
    p = tf.add_paragraph()
    p.text = '• 花生糖、芝麻糖、椰子糖'
    p.level = 1
    p = tf.add_paragraph()
    p.text = '• 甜蜜寓意，甜甜蜜蜜'
    p.level = 1
    p = tf.add_paragraph()
    p.text = '坚果'
    p.level = 1
    p = tf.add_paragraph()
    p.text = '• 瓜子、核桃、开心果'
    p.level = 2
    p = tf.add_paragraph()
    p.text = '• 寓意多子多福'
    p.level = 2
    p = tf.add_paragraph()
    p.text = '蜜饯'
    p.level = 1
    p = tf.add_paragraph()
    p.text = '• 传统零食，口味丰富'
    p.level = 2
    
    # Slide 6: 春节饮品类
    slide = prs.slides.add_slide(bullet_slide_layout)
    title_shape = slide.shapes.title
    body_shape = slide.placeholders[1]
    
    title_shape.text = '春节饮品类'
    
    tf = body_shape.text_frame
    tf.text = '茶类'
    p = tf.add_paragraph()
    p.text = '• 普洱茶、铁观音、龙井'
    p.level = 1
    p = tf.add_paragraph()
    p.text = '• 解腻助消化，待客首选'
    p.level = 1
    p = tf.add_paragraph()
    p.text = '酒类'
    p.level = 1
    p = tf.add_paragraph()
    p.text = '• 白酒、黄酒、红酒'
    p.level = 2
    p = tf.add_paragraph()
    p.text = '• 敬酒表达祝福'
    p.level = 2
    p = tf.add_paragraph()
    p.text = '汤品'
    p.level = 1
    p = tf.add_paragraph()
    p.text = '• 八宝鸭汤、老鸭汤'
    p.level = 2
    p = tf.add_paragraph()
    p.text = '• 滋补养生'
    p.level = 2
    
    # Slide 7: 春节饮食注意事项
    slide = prs.slides.add_slide(bullet_slide_layout)
    title_shape = slide.shapes.title
    body_shape = slide.placeholders[1]
    
    title_shape.text = '春节饮食注意事项'
    
    tf = body_shape.text_frame
    tf.text = '营养均衡'
    p = tf.add_paragraph()
    p.text = '• 荤素搭配，避免过度油腻'
    p.level = 1
    p = tf.add_paragraph()
    p.text = '• 多吃蔬菜水果'
    p.level = 1
    p = tf.add_paragraph()
    p.text = '适量原则'
    p.level = 1
    p = tf.add_paragraph()
    p.text = '• 控制食量，避免暴饮暴食'
    p.level = 2
    p = tf.add_paragraph()
    p.text = '• 少喝酒精饮料'
    p.level = 2
    p = tf.add_paragraph()
    p.text = '食品安全'
    p.level = 1
    p = tf.add_paragraph()
    p.text = '• 注意食材新鲜'
    p.level = 2
    p = tf.add_paragraph()
    p.text = '• 充分加热烹饪'
    p.level = 2
    
    # Slide 8: 春节饮食的地域差异
    slide = prs.slides.add_slide(bullet_slide_layout)
    title_shape = slide.shapes.title
    body_shape = slide.placeholders[1]
    
    title_shape.text = '春节饮食的地域差异'
    
    tf = body_shape.text_frame
    tf.text = '华北地区'
    p = tf.add_paragraph()
    p.text = '• 以饺子、年糕为主'
    p.level = 1
    p = tf.add_paragraph()
    p.text = '• 注重面食'
    p.level = 1
    p = tf.add_paragraph()
    p.text = '华南地区'
    p.level = 1
    p = tf.add_paragraph()
    p.text = '• 以汤圆、年糕为主'
    p.level = 2
    p = tf.add_paragraph()
    p.text = '• 重视海鲜'
    p.level = 2
    p = tf.add_paragraph()
    p.text = '华东地区'
    p.level = 1
    p = tf.add_paragraph()
    p.text = '• 菜系丰富多样'
    p.level = 2
    p = tf.add_paragraph()
    p.text = '• 精致美味'
    p.level = 2
    
    # Slide 9: 现代春节饮食新趋势
    slide = prs.slides.add_slide(bullet_slide_layout)
    title_shape = slide.shapes.title
    body_shape = slide.placeholders[1]
    
    title_shape.text = '现代春节饮食新趋势'
    
    tf = body_shape.text_frame
    tf.text = '健康饮食'
    p = tf.add_paragraph()
    p.text = '• 低油低盐，注重营养'
    p.level = 1
    p = tf.add_paragraph()
    p.text = '• 有机食材更受欢迎'
    p.level = 1
    p = tf.add_paragraph()
    p.text = '创新融合'
    p.level = 1
    p = tf.add_paragraph()
    p.text = '• 传统与现代结合'
    p.level = 2
    p = tf.add_paragraph()
    p.text = '• 中西合璧的新菜品'
    p.level = 2
    p = tf.add_paragraph()
    p.text = '便捷化'
    p.level = 1
    p = tf.add_paragraph()
    p.text = '• 半成品菜肴'
    p.level = 2
    p = tf.add_paragraph()
    p.text = '• 外卖年夜饭'
    p.level = 2
    
    # Slide 10: 结语
    end_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(end_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "春节美食，承载着千年文化"
    subtitle.text = "愿您的春节餐桌温馨圆满，年年有余！"
    
    # Customize title font
    title.text_frame.paragraphs[0].font.size = Pt(36)
    title.text_frame.paragraphs[0].font.bold = True
    
    # Save presentation
    file_path = '/Users/xuan.lx/Documents/x-agent/x-agent/workspace/chinese_new_year_food_presentation.pptx'
    prs.save(file_path)
    print(f"PPT已创建完成: {file_path}")

if __name__ == "__main__":
    create_chinese_new_year_food_presentation()